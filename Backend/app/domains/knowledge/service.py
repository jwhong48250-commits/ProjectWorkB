# app/domains/knowledge/service.py
"""Internal document ingestion helpers.

This module is kept small because the current router does not expose document
upload endpoints yet. The functions below are safe building blocks for the
future knowledge-base API and keep the merged code importable.
"""

from __future__ import annotations

from enum import nonmember
import io, re, os, subprocess, tempfile, asyncio
import json as _json
from datetime import datetime
from pathlib import Path
from typing import Any, Optional
from openai import AsyncOpenAI
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.core.config import settings
from app.domains.knowledge.agent_utils import chroma_client, get_collection
from app.utils.time_utils import now_kst
from app.domains.action.models import Priority

_async_openai = AsyncOpenAI(api_key=settings.OPENAI_API_KEY)

_splitter = RecursiveCharacterTextSplitter(
    chunk_size=800,
    chunk_overlap=100,
    separators=["\n\n", "\n", "。", ". ", " ", ""],
)


def _collection_name(workspace_id: int) -> str:
    """Return the workspace-scoped Chroma collection name."""
    return f"ws_{workspace_id}_docs"


def _extract_pdf(file_bytes: bytes) -> str:
    """Extract plain text from a PDF file."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_pptx(file_bytes: bytes) -> str:
    """
    PPT/PPTX -> 텍스트.

    슬라이드별 도형(shape) 텍스트 프레임 순회.
    "[슬라이드 N]" 헤더 삽입 -> 검색 결과에서 몇 번째 슬라이드인지 참조 가능.
    차트/이미지 속 텍스트는 추출 불가 -> vision 도메인 OCR 필요.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_html(file_bytes: bytes) -> str:
    """
    HTML -> 텍스트.

    제거 태그: script/style(코드 노이즈), nav/header/footer(반복 메뉴 텍스트).
    연속 줄바꿈 3개 이상 -> 2개로 압축해 빈 청크 방지.
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(file_bytes, "html.parser")
    for tag in soup(["script", "style", "nav", "header", "footer"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    return re.sub(rf"\n{3,}", "\n\n", text).strip()


def _extract_md(file_bytes: bytes) -> str:
    """
    Markdown -> 텍스트.

    별도 라이브러리 불필요 - UTF-8 텍스트 그대로 반환.
    마크다운 문법 기호(#, *, `, -)는 의미 있는 컨텍스트이므로 제거하지 않음.
    LLM이 마크다운을 이해하므로 오히려 구조 보존이 검색 품질에 유리.
    """
    return file_bytes.decode("utf-8", errors="replace")


def _extract_docx(file_bytes: bytes) -> str:
    """
    DOCX -> 텍스트.

    단락(paragraph) + 표(table) 두 가지 요소 순회.
    표는 행 단위로 셀을 탭으로 구분해 검색 가능한 텍스트로 변환.
    이미지/도형 안 텍스트는 추출 불가 -> vision 도메인 OCR 필요.

    주의: 구형 .doc 포맷은 python-docx 미지원 -> 422 반환
    """
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts = []

    # 단락 추출
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 표 추출 - 행 단위로 셀을 탭으로 이어붙임
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    return "\n\n".join(parts)


def _extrac_doc_legacy(file_bytes: bytes) -> str:
    """
    구형 .doc -> LibreOffice로 .docx 변환 후 python-docx로 추출.

    python-docx는 구현 .doc 포맷 미지원.
    LibreOffice --headless 모드로 .docx 변환 후 _extract_docx()에 위임.
    vision 도메인의 PPTX->PDF 변환과 동일한 LibreOffice 의존성.
    변환 타임아웃 30초 - 대용량 파일은 초과 가능.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        doc_path = os.path.join(tmpdir, "input, doc")
        with open(doc_path, "wb") as f:
            f.write(file_bytes)

        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "docx",
                "--outdir",
                tmpdir,
                doc_path,
            ],
            capture_output=True,
            timeout=30,
        )
        if result.returncode != 0:
            raise ValueError(
                "구형 .doc 변환 실패. 서버에 LibreOffice가 설치되어 있는지 확인하세요."
            )

        docx_path = os.path.join(tmpdir, "input-docx")
        with open(docx_path, "rb") as f:
            return _extract_docx(f.read())


def _extract_xlsx(file_bytes: bytes) -> str:
    """
    XLSX -> 텍스트.

    시트별로 "[시트명]" 헤더를 붙여 구분.
    수식 캐싱값 없는 파일은 LibreOffice로 재계산 후 재시도.
    LibreOffice 미설치 시 수식 문자열 fallback.
    """
    from openpyxl import load_workbook

    def _read(data: bytes) -> str:
        wb = load_workbook(io.BytesIO(data), data_only=True)
        sheets = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[{sheet.title}]\n" + "\n".join(rows))
        return "\n\n".join(sheets)

    text = _read(file_bytes)

    # None 셀 비율이 높으면 LibreOffice로 재계산 시도
    non_ratio = text.count("\t\t") / max(text.count("\t"), 1)
    if non_ratio > 0.3:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                xlsx_path = os.path.join(tmpdir, "input.xlsx")
                with open(xlsx_path, "wb") as f:
                    f.write(file_bytes)
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "xlsx",
                        "--outdir",
                        tmpdir,
                        xlsx_path,
                    ],
                    capture_output=True,
                    timeout=30,
                )
                out_path = os.path.join(tmpdir, "input.xlsx")
                with open(out_path, "rb") as f:
                    text = _read(f.read())
        except Exception:
            pass  # LibreOffice 없으면 원래 결과 그대로 사용

    return text


def _extract_xls_legacy(file_bytes: bytes) -> str:
    """
    구형 .xls -> LibreOffice로 .xlsx 변환 후 _extract_xlsx()에 위임

    xlrd 대신 LibreOffice를 쓰는 이유:
        1. xlrd는 수식 캐싱값을 못 읽는 케이스 동일하게 존재.
        2. LibreOffice 변환 시 수식 재계산까지 처리.
        3. xlrd 의존성 제거.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        xls_path = os.path.join(tmpdir, "input.xls")
        with open(xls_path, "wb") as f:
            f.write(file_bytes)

        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "xlsx",
                "--outdir",
                tmpdir,
                xls_path,
            ],
            capture_output=True,
            timeout=30,
        )
        if result.returncode != 0:
            raise ValueError(
                "구형 .xls 변환 실패. 서버에 LibreOffice가 설치되어 있는지 확인하세요."
            )

        xlsx_path = os.path.join(tmpdir, "input.xlsx")
        with open(xlsx_path, "rb") as f:
            return _extract_xlsx(f.read())


async def _embed_all(texts: list[str]) -> list[list[float]]:
    """
    청크 전체를 OpenAI 임베딩 API로 병렬 배치 처리.
    OpenAI 단일 요청 최대 2048개 제한에 맞게 분할 후 asyncio.gather로 동시 실행.
    ChromaDB EF의 순차 처리 대비 10배 이상 빠름
    """
    BATCH = 2048
    batches = [texts[i : i + BATCH] for i in range(0, len(texts), BATCH)]
    responses = await asyncio.gather(
        *[
            _async_openai.embeddings.create(
                model="text-embedding-3-small",
                input=batch,
            )
            for batch in batches
        ]
    )
    return [item.embedding for resp in responses for item in resp.data]


async def ingest_document(
    workspace_id: int,
    filename: str,
    file_bytes: bytes,
    file_type: str,
    title: Optional[str] = None,
) -> dict:
    """
    문서 수집 전체 파이프라인.

    Args:
        workspace_id: 워크스페이스 ID, 컬렉션 격리 키.
        filename: 웝본 파일명. doc_id 및 메타데이터에 사용.
        file_bytes: 업로드된 파일 바이너리.
        file_type: "pdf" | "pptx" | "html".
        title: 문서 제목. None 이면 filename 사용.

    Returns:
        {"doc_id": str, "chunks": int, "title": str}

    중복 업로드 처리:
        doc_id = "{workspace_id}_{filename}" 고정
        청크 ID도 "{doc_id}_{chunk_id}" 고정.
        -> 같은 파일 재업로드 시 upsert가 덮어씀. 벡터 중복 없음
    """
    # 1단계: 텍스트 추출
    extractors = {
        "pdf": _extract_pdf,
        "pptx": _extract_pptx,
        "html": _extract_html,
        "md": _extract_md,
        "docx": _extract_docx,
        "doc": _extrac_doc_legacy,
        "xlsx": _extract_xlsx,
        "xls": _extract_xls_legacy,
    }
    extractor = extractors.get(file_type)
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type}")
    raw_text = extractor(file_bytes)

    if not raw_text.strip():
        raise ValueError("텍스트를 추출할 수 없습니다.")

    # 2단계: 청크 분할
    # 결과: 각 청크 <= 800자, 인접 청크 간 100자 overlap
    chunks = _splitter.split_text(raw_text)

    # 3단계: 임베딩 - async 배치 병렬 처리
    embeddings = await _embed_all(chunks)

    # 4단계: 메타데이터 구성
    doc_id = f"{workspace_id}_{filename}"
    title = title or filename
    uploaded_at = now_kst().isoformat()

    ids = [f"{doc_id}_chunk{i}" for i in range(len(chunks))]
    metadatas = [
        {
            "workspace_id": workspace_id,
            "doc_id": doc_id,
            "title": title,
            "filename": filename,
            "file_type": file_type,
            "source_type": "uploaded",
            "chunk_index": i,
            "total_chunks": len(chunks),
            "uploaded_at": uploaded_at,
        }
        for i in range(len(chunks))
    ]

    # 5단계: ChromaDB upsert
    # add() 대신 upsert() 이유:
    #   add()는 동일 ID 존재 시 에러 → 재업로드 불가
    #   upsert()는 있으면 덮어쓰고 없으면 삽입 → 재업로드 안전
    collection = get_collection(workspace_id)
    collection.upsert(
        documents=chunks, embeddings=embeddings, ids=ids, metadatas=metadatas
    )

    return {"doc_id": doc_id, "chunks": len(chunks), "title": title}


def _extract_wbs_json(content: str) -> str:
    """
    WBS JSON -> 검색 가능한 평문.
    {"epics": [{"title": "...", "tasks": [{"title": "...", "assignee": "..."}]}]}
    """
    try:
        wbs = _json.loads(content)
    except Exception:
        return content  # JSON 파싱 실패 시 raw 반환
    lines = []
    for epic in wbs.get("epics", []):
        lines.append(f"[에픽] {epic.get("title", '')}")
        for task in epic.get("tasks", []):
            assignee = task.get("assignee", "")
            title = task.get("title", "")
            priority = task.get("priority", "")
            deadline = task.get("due_date", "")
            lines.append(
                f" - {title} / 담당: {assignee} / 우선순위: {priority} / 마감: {deadline}"
            )
    return "\n".join(lines)


async def ingest_db_content(
    workspace_id: int,
    doc_id: str,  # "minutes_{meeting_id}" | "report_{report_id}"
    content: str,
    format: str,  # "markdown" | "html" | "wbs"  (excel은 content 없으므로 호출 안 함)
    title: str = "",
    extra_metadata: dict = {},
) -> dict:
    """
    content -> ChromaDB 인제스트.
    meeting_minutes, reports 공통 파이프라인.
    format별로 텍스트 추출 방법이 다름:
        markdown -> _extract_md
        html -> _extract_html
        wbs -> _extract_wbs_json (JSON -> 평문)
    """
    if format == "markdown":
        raw_text = _extract_md(content.encode("utf-8"))
    elif format == "html":
        raw_text = _extract_html(content.encode("utf-8"))
    elif format == "wbs":
        raw_text = _extract_wbs_json(content)
    else:
        raise ValueError(f"인제스트 미지원 포맷: {format}")

    if not raw_text.strip():
        raise ValueError("보고서 텍스트가 비어있습니다.")

    chunks = _splitter.split_text(raw_text)
    embeddings = await _embed_all(chunks)
    uploaded_at = now_kst().isoformat()

    ids = [f"{doc_id}_chunk{i}" for i in range(len(chunks))]
    metadatas = [
        {
            "workspace_id": workspace_id,
            "doc_id": doc_id,
            "title": title,
            "chunk_index": i,
            "total_chunks": len(chunks),
            "uploaded_at": uploaded_at,
            **extra_metadata,
        }
        for i in range(len(chunks))
    ]

    collection = get_collection(workspace_id)
    collection.upsert(
        documents=chunks, embeddings=embeddings, ids=ids, metadatas=metadatas
    )
    return {"doc_id": doc_id, "chunks": len(chunks)}


async def analyze_document_for_display(
    workspace_id: int,
    filename: str,
    file_bytes: bytes,
    file_type: str,
    title: str | None = None,
) -> dict:
    """
    텍스트 추출 + LLM 요약 -> 화면 표시용 분석 결과.
    200,000자 이하: 전체 전송 (gpt-4o-mini 128k 컨텍스트 내)
    초과 시: Map-Reduce (청크별 미니 요약 → 최종 요약)
    """
    extractors = {
        "pdf": _extract_pdf,
        "pptx": _extract_pptx,
        "html": _extract_html,
        "md": _extract_md,
        "docx": _extract_docx,
        "doc": _extrac_doc_legacy,
        "xlsx": _extract_xlsx,
        "xls": _extract_xls_legacy,
    }.get(file_type)
    if not extractors:
        raise ValueError(f"Unsupported file type: {file_type}")

    raw_text = extractors(file_bytes)
    if not raw_text.strip():
        raise ValueError("텍스트를 추출할 수 없습니다.")

    MAX_DIRECT = 200_000

    if len(raw_text) <= MAX_DIRECT:
        # 전체 텍스트 직접 전송
        content_for_llm = raw_text
    else:
        # Map-Reduce: 청크별 미니 요약 -> 합치기
        chunks = _splitter.split_text(raw_text)
        mini_responses = await asyncio.gather(
            *[
                _async_openai.chat.completions.create(
                    model="gpt-5.4-mini",
                    messages=[
                        {
                            "role": "user",
                            "content": f"다음 내용에서 중요한 정보를 빠짐없이 보존하며 핵심만 간결하게 정리하세요:\n{chunk}",
                        }
                    ],
                )
                for chunk in chunks
            ]
        )
        content_for_llm = "\n\n".join(
            r.choices[0].message.content for r in mini_responses
        )

    result = await _async_openai.chat.completions.create(
        model="gpt-5.4-mini",
        messages=[
            {
                "role": "user",
                "content": f"""
            다음은 문서 전체를 섹션별로 정리한 내용이다. 전체 문서를 바탕으로 분석해 JSON으로 답하세요.

            문서용: {title or filename}
            내용:
            {content_for_llm}

            {{
                "summary": "문서 전체의 목적과 주요 내용을 5-7문장으로 요약",
                "key_points": ["핵심 포인트 (최대 7개, 구체적으로)"]
            }}
            """,
            }
        ],
        response_format={"type": "json_object"},
    )

    try:
        parsed = _json.loads(result.choices[0].message.content)
    except Exception:
        parsed = {"summary": result.choices[0].message.content, "key_points": []}

    return {
        "filename": filename,
        "title": title or filename,
        "summary": parsed.get("summary", ""),
        "key_points": parsed.get("key_points", []),
    }


async def process_meeting_end(meeting_id: int, workspace_id: int) -> None:
    """
    회의 종료 후리 (BackgroundTask):
    utterances(MongoDB) → LLM → decisions + wbs_tasks + summary → MySQL 저장.
    회의 종료 버튼 한 번으로 모든 데이터가 DB에 적재된다.
    """
    from motor.motor_asyncio import AsyncIOMotorClient
    from langchain_openai import ChatOpenAI
    from sqlalchemy.orm import Session
    import json, re
    from datetime import date

    from app.core.config import settings
    from app.infra.database.session import SessionLocal
    from app.domains.intelligence.models import Decision, MeetingMinute, MinuteStatus
    from app.domains.action.models import WbsEpic, WbsTask, Priority
    from app.utils.time_utils import now_kst

    # Phase 1: MongoDB fetch — no DB session held
    mongo_client = AsyncIOMotorClient(settings.MONGODB_URL)
    try:
        mongo_db = mongo_client["meeting_assistant"]
        ctx_doc = await mongo_db["utterances"].find_one(
            {"$or": [{"meeting_id": meeting_id}, {"meeting_id": str(meeting_id)}]}
        )
    finally:
        mongo_client.close()

    if not ctx_doc or not ctx_doc.get("utterances"):
        return

    transcript_text = "\n".join(
        f"[{u.get('speaker_label', '?')}] {u.get('content', '')}"
        for u in ctx_doc["utterances"]
    )

    # Phase 2: LLM call — no DB session held
    llm = ChatOpenAI(model="gpt-4o-mini", api_key=settings.OPENAI_API_KEY)
    today_str = now_kst().strftime("%Y-%m-%d")
    prompt = f"""
    다음 회의 발화에서 구조화된 정보를 추출하세요.

    오늘 날짜: {today_str}
    due_date 계산 시 이 날짜를 기준으로 상대적 표현("다음 주까지", "3일 안에" 등)을 절대 날짜로 변환하세요.

    [발화 내용]
    {transcript_text}

    액션 아이템 우선순위(priority) 판단 기준:
    - high: 결정 사항과 직접 연결 / 다른 액션의 선행 조건 / "반드시·꼭·최우선" 발화 / 다수 인원 영향
    - normal: 그 외

    긴급도(urgency) 판단 기준:
    - urgent: 기한 3일 이내 / 다음 회의 전 완료 필요 / "빨리·즉시·오늘까지·ASAP·as soon as possible" 발화
    - normal: 기한 4~7일 이내
    - low: 기한 7일 초과 또는 미언급

    반드시 아래 JSON 형식으로만 답변하세요.
    {{
        "title": "회의 핵심 주제 (한 줄)",
        "key_points": ["핵심 논의 내용 1", ...],
        "decisions": [
            {{"content": "결정 사항 내용"}}
        ],
        "wbs_tasks": [
            {{
                "order": 0,
                "title": "할 일 제목 (한 줄, 30자 이내)",
                "content": "할 일 상세 내용 (담당자·기한·배경 포함, 없으면 null)",
                "assignee_name": "담당자 이름 or null",
                "due_date": "YYYY-MM-DD or null",
                "priority": "low|medium|high|critical",
                "urgency": "urgent|normal|low"
            }}
        ],
        "hallucination_flags": ["근거 불충분 항목 설명 (없으면 빈 배열)"]
    }}
    """
    result = await llm.ainvoke(prompt)
    json_match = re.search(r"\{.*\}", result.content, re.DOTALL)
    try:
        extracted = json.loads(json_match.group()) if json_match else {}
    except json.JSONDecodeError:
        extracted = {}

    # Phase 3: DB writes — session opened only for the write phase
    db = SessionLocal()
    try:
        now = now_kst().replace(tzinfo=None)

        # -- decisions -> MySQL --
        for d in extracted.get("decisions", []):
            db.add(
                Decision(
                    meeting_id=meeting_id,
                    content=d["content"],
                    speaker_id=None,
                    detected_at=now,
                    is_confirmed=False,
                )
            )

        # -- wbs_tasks -> wbs_epics _ wbs_tasks -> MySQL
        wbs_tasks = extracted.get("wbs_tasks", [])
        if wbs_tasks:
            epic = WbsEpic(
                meeting_id=meeting_id,
                title="회의 액션 아이템",
                order_index=0,
            )
            db.add(epic)
            db.flush()

            for item in wbs_tasks:
                due = None
                if item.get("due_date"):
                    try:
                        due = date.fromisoformat(item["due_date"])
                    except Exception:
                        pass

                # priority는 Priority enum으로 변환, 범위 밖 값이면 medium 기본값
                try:
                    priority = Priority(item.get("priority", "medium"))
                except ValueError:
                    priority = Priority.medium

                db.add(
                    WbsTask(
                        epic_id=epic.id,
                        title=item.get("title") or item.get("content", "")[:200],
                        content=item.get("content"),
                        assignee_name=item.get("assignee_name"),
                        priority=priority,
                        urgency=item.get("urgency", "normal"),
                        due_date=due,
                        order_index=item.get("order", 0),
                    )
                )

        # summary만 meeting_minutes에 저장 (decisions/wbs_tasks는 별도 테이블에 있으므로 제외)
        summary_only = {
            "title": extracted.get("title", ""),
            "key_points": extracted.get("key_points", []),
            "hallucination_flags": extracted.get("hallucination_flags", []),
        }
        minute = (
            db.query(MeetingMinute)
            .filter(MeetingMinute.meeting_id == meeting_id)
            .one_or_none()
        )
        if minute is None:
            minute = MeetingMinute(meeting_id=meeting_id, status=MinuteStatus.draft)
            db.add(minute)
        minute.summary = json.dumps(summary_only, ensure_ascii=False)

        db.commit()

    except Exception:
        db.rollback()
        raise
    finally:
        db.close()
